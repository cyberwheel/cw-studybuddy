# --- indexer.py (Final Version) ---

import os
import csv
from pathlib import Path
import pdfplumber
import docx
import pptx
import chromadb
from sentence_transformers import SentenceTransformer

# --- Setup ---
DB_PATH = "db"
NOTES_PATH = Path("My_Notes")
QA_FILE = NOTES_PATH / "qa_pairs.csv"

print("📘 Initializing database and model...")
db_client = chromadb.PersistentClient(path=DB_PATH)
embedding_model = SentenceTransformer("all-MiniLM-L6-v2")
print("✅ Model loaded successfully.")

# --- Helper: Text Splitter ---
def split_text_into_chunks(text, chunk_size=1000, overlap=100):
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += (chunk_size - overlap)
    return chunks

# --- Helper: Read Q&A CSV ---
def read_qa_pairs():
    qas = []
    if QA_FILE.exists():
        try:
            with open(QA_FILE, "r", encoding="utf-8") as f:
                reader = csv.reader(f)
                next(reader, None)  # Skip header
                for row in reader:
                    if len(row) >= 2:
                        qas.append((row[0].strip(), row[1].strip()))
        except Exception as e:
            print(f"⚠️ Error reading Q&A file: {e}")
    return qas

# --- Main Processor ---
def process_materials():
    NOTES_PATH.mkdir(exist_ok=True)
    print("📂 Scanning My_Notes...")

    doc_texts, doc_meta = [], []

    # Read PDFs
    for f in NOTES_PATH.glob("*.pdf"):
        try:
            with pdfplumber.open(f) as pdf:
                text = "\n".join([p.extract_text() or "" for p in pdf.pages])
            doc_texts.append(text)
            doc_meta.append({"source": f.name})
            print(f"📄 Loaded: {f.name}")
        except Exception as e:
            print(f"❌ {f.name}: {e}")

    # Read Word docs
    for f in NOTES_PATH.glob("*.docx"):
        try:
            d = docx.Document(f)
            text = "\n".join([p.text for p in d.paragraphs])
            doc_texts.append(text)
            doc_meta.append({"source": f.name})
            print(f"📝 Loaded: {f.name}")
        except Exception as e:
            print(f"❌ {f.name}: {e}")

    # Read PowerPoints
    for f in NOTES_PATH.glob("*.pptx"):
        try:
            pres = pptx.Presentation(f)
            text = "\n".join(
                [shape.text for slide in pres.slides for shape in slide.shapes if hasattr(shape, "text")]
            )
            doc_texts.append(text)
            doc_meta.append({"source": f.name})
            print(f"🎞️ Loaded: {f.name}")
        except Exception as e:
            print(f"❌ {f.name}: {e}")

    # Read Q&A pairs
    qas = read_qa_pairs()
    for q, a in qas:
        doc_texts.append(f"Question: {q}\nAnswer: {a}")
        doc_meta.append({"source": "qa_pairs.csv"})
    print(f"💬 Loaded {len(qas)} Q&A pairs.")

    # Split, embed, and store
    if not doc_texts:
        print("⚠️ No valid documents found.")
        return

    all_chunks, all_meta, all_ids = [], [], []
    for idx, (text, meta) in enumerate(zip(doc_texts, doc_meta)):
        for j, chunk in enumerate(split_text_into_chunks(text)):
            if chunk.strip():
                all_chunks.append(chunk)
                all_meta.append(meta)
                all_ids.append(f"{meta['source']}_{idx}_{j}")

    print(f"🔍 Creating {len(all_chunks)} embeddings...")
    embeds = embedding_model.encode(all_chunks, batch_size=64, show_progress_bar=True).tolist()

    try:
        db_client.delete_collection("course_notes")
    except:
        pass

    coll = db_client.create_collection("course_notes", metadata={"hnsw:space": "cosine"})
    for i in range(0, len(all_chunks), 100):
        coll.add(
            documents=all_chunks[i:i+100],
            metadatas=all_meta[i:i+100],
            ids=all_ids[i:i+100],
            embeddings=embeds[i:i+100],
        )

    print(f"✅ Successfully stored {len(all_chunks)} chunks in vector DB.")

if __name__ == "__main__":
    process_materials()
